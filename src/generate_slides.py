import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 Aspect Ratio dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Color Palette
COLOR_BG_DARK = RGBColor(15, 23, 42)       # Slate 900
COLOR_BG_LIGHT = RGBColor(248, 250, 252)   # Slate 50
COLOR_CARD_BG = RGBColor(255, 255, 255)    # White
COLOR_CARD_DARK = RGBColor(30, 41, 59)     # Slate 800
COLOR_PRIMARY = RGBColor(37, 99, 235)      # Blue 600
COLOR_PRIMARY_DARK = RGBColor(29, 78, 216) # Blue 700
COLOR_ACCENT = RGBColor(16, 185, 129)      # Emerald 500
COLOR_WARNING = RGBColor(245, 158, 11)     # Amber 500
COLOR_TEXT_DARK = RGBColor(15, 23, 42)     # Slate 900
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
COLOR_TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100
COLOR_BORDER = RGBColor(226, 232, 240)     # Slate 200

def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_header(slide, title_text, category_text):
        # Category pill/badge
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(239, 246, 255)
        badge.line.color.rgb = COLOR_PRIMARY
        badge.line.width = Pt(1)
        tf_badge = badge.text_frame
        tf_badge.word_wrap = True
        p_badge = tf_badge.paragraphs[0]
        p_badge.text = category_text.upper()
        p_badge.font.size = Pt(10)
        p_badge.font.bold = True
        p_badge.font.color.rgb = COLOR_PRIMARY
        p_badge.font.name = "Segoe UI"
        p_badge.alignment = PP_ALIGN.CENTER

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_DARK
        p.font.name = "Segoe UI"

    # =========================================================================
    # SLIDE 1: Title Slide (Dark Theme)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG_DARK
    bg1.line.fill.background()

    # Decorative accent line
    dec = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(0.8), Inches(0.08))
    dec.fill.solid()
    dec.fill.fore_color.rgb = COLOR_PRIMARY
    dec.line.fill.background()

    # Title & Subtitle Box
    tbox = s1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.2))
    tf1 = tbox.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "Dogs & Cats Breed Classifier"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_LIGHT
    p1.font.name = "Segoe UI"

    p2 = tf1.add_paragraph()
    p2.text = "End-to-End Machine Learning Data Pipeline & Strategic Next-Phase Roadmap"
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_PRIMARY
    p2.font.name = "Segoe UI"
    p2.space_before = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "สรุปภาพรวม Pipeline • การแก้ไขปัญหาและอุปสรรค • แนวทางต่อยอดโมเดล • แผนการปรับปรุงในอนาคต"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.font.name = "Segoe UI"
    p3.space_before = Pt(8)

    # Info footer card
    info_card = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.6), Inches(11.33), Inches(1.1))
    info_card.fill.solid()
    info_card.fill.fore_color.rgb = COLOR_CARD_DARK
    info_card.line.color.rgb = RGBColor(51, 65, 85)
    info_card.line.width = Pt(1)
    tf_info = info_card.text_frame
    p_info = tf_info.paragraphs[0]
    p_info.text = "Dataset: rajarshi2712/dogs-and-cats-classifier (28,983 Images | 104 Classes)  •  Automation: run_pipeline.bat\nGitHub: github.com/wayu0849-lang/-  •  Date: August - September 2026"
    p_info.font.size = Pt(12)
    p_info.font.color.rgb = RGBColor(203, 213, 225)
    p_info.font.name = "Segoe UI"
    p_info.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: Pipeline Overview (Key Architecture & Statistics)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "ภาพรวมของระบบ Data Pipeline (Architecture & Highlights)", "Pipeline Overview")

    steps = [
        ("1. Data Collection", "• ดึงข้อมูลผ่าน kagglehub API\n• จัดเก็บสู่เครื่องท้องถิ่นใน data/\n• กรอง Git ด้วย .gitignore ป้องกัน Push รูปดิบ", COLOR_PRIMARY),
        ("2. Exploratory EDA", "• ตรวจสอบ 28,983 รูป (0 ไฟล์เสีย)\n• วิเคราะห์ 104 คลาส สุนัข/แมว\n• วัด Class Imbalance (12.41x)\n• สรุปขนาด 224x224 RGB 100%", COLOR_PRIMARY_DARK),
        ("3. Preprocessing", "• มาตรฐาน RGB 3 ช่องสี\n• Resize Lanczos & Aspect-Pad\n• Scale [0, 1] + ImageNet Norm\n• Data Augmentation (357 FPS)", COLOR_ACCENT),
        ("4. Data Splitting", "• Train 92.3% (26,752 ภาพ)\n• Valid 4.6% (1,339 ภาพ)\n• Test 3.1% (892 ภาพ)\n• 100% Stratification & 0 Leak", RGBColor(124, 58, 237)),
    ]

    card_w = Inches(2.75)
    card_h = Inches(3.8)
    gap = Inches(0.24)
    start_x = Inches(0.8)

    for i, (title, content, border_c) in enumerate(steps):
        cx = start_x + i * (card_w + gap)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, Inches(1.8), card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = border_c
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.bold = True
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = border_c
        p_t.font.name = "Segoe UI"

        p_c = tf.add_paragraph()
        p_c.text = content
        p_c.font.size = Pt(11)
        p_c.font.color.rgb = COLOR_TEXT_DARK
        p_c.font.name = "Segoe UI"
        p_c.space_before = Pt(10)

    # Bottom summary metric banner
    banner = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.9), Inches(11.73), Inches(1.0))
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(241, 245, 249)
    banner.line.color.rgb = COLOR_BORDER
    tf_b = banner.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "⚡ ประสิทธิภาพและระบบอัตโนมัติ: มี run_pipeline.bat สั่งรันรวดเดียวแบบ End-to-End ครบทุกขั้นตอน | ผ่านการทดสอบ Throughput 357 ภาพ/วินาทีบน CPU"
    p_b.font.size = Pt(11.5)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_TEXT_DARK
    p_b.font.name = "Segoe UI"
    p_b.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: Detailed EDA & Preprocessing Visuals
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "ข้อค้นพบจากการสำรวจข้อมูล และการเตรียมรูปภาพ (EDA & Preprocessing)", "Technical Deep Dive")

    # Card Left: EDA Findings
    c_left = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.2))
    c_left.fill.solid()
    c_left.fill.fore_color.rgb = COLOR_CARD_BG
    c_left.line.color.rgb = COLOR_BORDER
    tf_l = c_left.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "📊 ข้อค้นพบสำคัญจาก EDA (Exploratory Data Analysis)"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Segoe UI"

    bullets_l = (
        "• ความสมบูรณ์ 100%: ไม่พบภาพชำรุดเสียหาย (0 Corrupted Images)\n"
        "• สัดส่วนสายพันธุ์: สุนัข 67.8% (19,658 ภาพ) | แมว 32.2% (9,325 ภาพ)\n"
        "• อัตราส่วน Class Imbalance: สูงสุด 12.41 เท่า\n"
        "  - คลาสสูงสุด: cat - bengal (1,080 ภาพ)\n"
        "  - คลาสต้อยต่ำสุด: dog - bernard-dog - saint (87 ภาพ)\n"
        "• ขนาดภาพมาตรฐาน: 224 x 224 พิกเซล, รูปแบบสี RGB 3 ช่องสี\n"
        "• สรุป: ต้องทำ Data Augmentation ในช่วง Train เพื่อป้องกัน Overfitting จากคลาสไม่สมดุล"
    )
    p_b = tf_l.add_paragraph()
    p_b.text = bullets_l
    p_b.font.size = Pt(11)
    p_b.font.color.rgb = COLOR_TEXT_DARK
    p_b.font.name = "Segoe UI"
    p_b.space_before = Pt(8)

    # Embed species breakdown or split distribution if exists
    img_eda = Path("reports/eda/species_breakdown.png")
    if img_eda.exists():
        s3.shapes.add_picture(str(img_eda.resolve()), Inches(1.8), Inches(4.5), width=Inches(3.6))

    # Card Right: Preprocessing & Augmentation
    c_right = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2))
    c_right.fill.solid()
    c_right.fill.fore_color.rgb = COLOR_CARD_BG
    c_right.line.color.rgb = COLOR_BORDER
    tf_r = c_right.text_frame
    tf_r.word_wrap = True

    p_r0 = tf_r.paragraphs[0]
    p_r0.text = "⚙️ กระบวนการ Preprocessing & Augmentation Pipeline"
    p_r0.font.bold = True
    p_r0.font.size = Pt(13)
    p_r0.font.color.rgb = COLOR_ACCENT
    p_r0.font.name = "Segoe UI"

    bullets_r = (
        "• Color Standard: บังคับแปลงเป็น RGB ป้องกันข้อผิดพลาด Grayscale/RGBA\n"
        "• Scaling & Normalization: แปลงค่าพิกเซล [0, 255] -> [0, 1] และ Z-Score ImageNet Standard (Mean=[0.485, 0.456, 0.406], Std=[0.229, 0.224, 0.225])\n"
        "• Data Augmentation (Training Pipeline):\n"
        "  - Random Horizontal Flip (p=0.5)\n"
        "  - Random Rotation (+/- 15 องศา)\n"
        "  - Color Jitter (Brightness/Contrast 0.85 - 1.15)\n"
        "• ความเร็วการประมวลผล: ~357 FPS พร้อมใช้งานใน DataLoader ทันที"
    )
    p_r1 = tf_r.add_paragraph()
    p_r1.text = bullets_r
    p_r1.font.size = Pt(11)
    p_r1.font.color.rgb = COLOR_TEXT_DARK
    p_r1.font.name = "Segoe UI"
    p_r1.space_before = Pt(8)

    img_prep = Path("reports/preprocessing/augmentation_samples.png")
    if img_prep.exists():
        s3.shapes.add_picture(str(img_prep.resolve()), Inches(7.0), Inches(4.8), width=Inches(5.3))

    # =========================================================================
    # SLIDE 4: Problems & Obstacles Encountered (ข้อ 1 ในโจทย์)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "ปัญหาและอุปสรรคที่พบระหว่างทำงาน และแนวทางแก้ไข (Challenges & Solutions)", "Problems & Obstacles")

    problems = [
        ("1. สิทธิ์ Administrator ในการติดตั้ง Git",
         "อุปสรรค: winget เรียก InnoSetup installer ของ Git ที่ร้องขอสิทธิ์ Admin (UAC Prompt) ทำให้คำสั่งล้มเหลวแบบ Non-interactive",
         "วิธีแก้ไข: ดาวน์โหลด MinGit (Portable Git 64-bit) มาแตกไฟล์ใน %LOCALAPPDATA%\\Programs\\Git และตั้งค่า User PATH ถาวร ทำให้ใช้งาน Git ได้สมบูรณ์โดยไม่ต้องพึ่งสิทธิ์แอดมิน"),

        ("2. ปัญหา Git Authentication & User Config",
         "อุปสรรค: การรัน git push ค้างเนื่องจาก Git ยังไม่ได้ตั้งค่า identity (Author unknown) และขาด Credential Helper สำหรับ HTTPS",
         "วิธีแก้ไข: กำหนด user.name 'wayu0849-lang' และ user.email ให้ตรงกับ GitHub บัญชีเดียวกัน พร้อมเปิดใช้งาน Git Credential Manager (manager)"),

        ("3. การป้องกันข้อมูลรูปภาพขนาดใหญ่หลุดขึ้น Git",
         "อุปสรรค: รูปภาพ 28,983 รูปขนาดรวมหลายร้อย MB หาก push ขึ้น GitHub จะเกินขีดจำกัดของ Git Repository",
         "วิธีแก้ไข: ออกแบบ .gitignore กรองโฟลเดอร์ data/ และนามสกุลรูปภาพ แต่ทำการ Whitelist เฉพาะไฟล์รายงานและชาร์ตใน reports/ ให้ติดตามได้ครบถ้วน"),

        ("4. ปัญหา Class Imbalance ในชุดข้อมูล (12.41x)",
         "อุปสรรค: สายพันธุ์สุนัขบางคลาสมีภาพเพียง 87 ภาพ ขณะที่แมวบางคลาสมีถึง 1,080 ภาพ เสี่ยงต่อการเกิดโมเดล Biased",
         "วิธีแก้ไข: ออกแบบ Stochastic Data Augmentation (Rotation, Jitter, Flip) ใน Training Pipeline และเสนอแนะการใช้ Focal Loss / Weighted Sampler"),

        ("5. ความเปลี่ยนแปลงของ API ใน Pandas 3.0",
         "อุปสรรค: การใช้ groupby().apply() ในการสุ่มภาพตรวจสอบ Data Leakage เกิดข้อผิดพลาดของ Multi-index ในเวอร์ชันล่าสุด",
         "วิธีแก้ไข: ปรับปรุงตรรกะการสุ่มตัวอย่างด้วย Dataframe Slicing ตรงตามราย Split ทำให้การตรวจสอบ MD5 Hash ราบรื่นและแม่นยำ")
    ]

    card_y = Inches(1.6)
    c_h = Inches(0.95)
    c_w = Inches(11.73)
    gap_y = Inches(0.14)

    for i, (p_title, p_desc, p_sol) in enumerate(problems):
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), card_y + i * (c_h + gap_y), c_w, c_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = p_title
        p1.font.bold = True
        p1.font.size = Pt(11.5)
        p1.font.color.rgb = COLOR_PRIMARY
        p1.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = p_desc + "  ➜  " + p_sol
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(3)

    # =========================================================================
    # SLIDE 5: Next-Phase Project Extensions (ข้อ 2 ในโจทย์ - แนวทางต่อยอด)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "แนวทางการนำชุดข้อมูลที่ Clean แล้วไปต่อยอดทำโปรเจกต์ (Next-Phase Project Roadmap)", "Project Extensions")

    projects = [
        ("Project 1: Fine-Tuning SOTA Vision Models (Transfer Learning)",
         "การสร้างโมเดลจำแนกสายพันธุ์ด้วย Pretrained Vision Backbones เช่น EfficientNetV2, ConvNeXt, Swin Transformer หรือ MobileNetV4",
         "• ความสอดคล้อง: ภาพขนาด 224x224 RGB ตรงกับ input architecture ของโมเดล SOTA อยู่แล้ว\n• ประสิทธิภาพ: Transfer Learning ช่วยให้โมเดลจำแนก 104 คลาสได้แม่นยำสูง (>90%) แม้ในคลาสที่มีรูปภาพน้อย (87 รูป)\n• ความคุ้มค่า: ประหยัดเวลาและพลังงาน Compute ได้มากกว่าเทรนจากศูนย์หลายเท่าตัว",
         COLOR_PRIMARY),

        ("Project 2: Hierarchical / Multi-Task Learning",
         "การจำแนกแบบลำดับชั้น (Species -> Breed -> Visual Attributes) เช่น จำแนกชนิดสัตว์, สายพันธุ์, สีขน, และลวดลาย",
         "• สกัด Feature ร่วม: พัฒนา Multi-Head Classifier เรียนรู้ทั้งหมวดหมู่หลัก (หมา/แมว) และหมวดหมู่ย่อยพร้อมกัน\n• Hierarchical Fallback: หากโมเดลไม่มั่นใจในระดับสายพันธุ์ ก็ยังสามารถทำนายชนิดสัตว์และลักษณะภายนอกได้ถูกต้อง\n• Robustness: ลดความสับสนระหว่างสายพันธุ์ที่มีรูปลักษณ์ใกล้เคียงกัน เช่น Siberian Husky กับ Alaskan Malamute",
         COLOR_ACCENT),

        ("Project 3: Real-Time Mobile & Edge Application (Edge AI)",
         "การแปลงโมเดลเป็น ONNX Runtime / TensorFlow Lite เพื่อพัฒนาแอปพลิเคชันมือถือหรือ Web App สำหรับระบุสายพันธุ์สัตว์เลี้ยง",
         "• Throughput สูง: Preprocessing ทำงานได้เร็วถึง 357 FPS บน CPU ทำให้รันบนอุปกรณ์ปลายทางได้แบบ Real-time\n• On-Device Offline: โมเดลอย่าง MobileNet มีขนาดเล็ก (<20 MB) สามารถสแกนและระบุสายพันธุ์ได้โดยไม่ต้องต่อเน็ต\n• Real-world Value: ใช้งานได้จริงในคลินิกสัตว์, โรงพยาบาลสัตว์, คาเฟ่สัตว์เลี้ยง, และศูนย์รับเลี้ยงสัตว์",
         RGBColor(124, 58, 237)),

        ("Project 4: Explainable AI (XAI) & Pet Facial Landmark Analysis",
         "การประยุกต์ใช้ Grad-CAM และ Integrated Gradients เพื่อตรวจสอบจุดสังเกต (Attention Heatmap) ที่โมเดลใช้ในการตัดสินใจ",
         "• ความโปร่งใส (Transparency): ตรวจสอบว่าโมเดลดูที่โครงหน้า, ทรงหู, หรือลายขนจริง แทนที่จะจำ Background หรือสิ่งแปลกปลอม\n• Medical Screening: สามารถต่อยอดวิเคราะห์โครงสร้างใบหน้าสัตว์เลี้ยงเพื่อคัดกรองปัญหาสุขภาพเบื้องต้นได้",
         RGBColor(217, 119, 6))
    ]

    card_w5 = Inches(5.7)
    card_h5 = Inches(2.4)
    positions5 = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.45)),
        (Inches(6.8), Inches(4.45))
    ]

    for (title, desc, details, color_theme), (px, py) in zip(projects, positions5):
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, py, card_w5, card_h5)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = color_theme
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(11.5)
        p1.font.color.rgb = color_theme
        p1.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.bold = True
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(3)

        p3 = tf.add_paragraph()
        p3.text = details
        p3.font.size = Pt(9)
        p3.font.color.rgb = COLOR_TEXT_MUTED
        p3.font.name = "Segoe UI"
        p3.space_before = Pt(3)

    # =========================================================================
    # SLIDE 6: Supporting Rationales (เหตุผลสนับสนุนการต่อยอด)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "เหตุผลสนับสนุนในการเลือกต่อยอดทางเทคนิคและธุรกิจ (Supporting Rationales)", "Supporting Rationales")

    pillars = [
        ("ความพร้อมทางเทคนิค (Technical Readiness)",
         "• ข้อมูลสะอาด 100%: ไม่มีไฟล์เสีย (Zero Corrupt) และมีมิติ 224x224 RGB สม่ำเสมอ\n• สถาปัตยกรรมรองรับ: มีฟังก์ชัน Preprocessing & Normalization ที่คำนวณสถิติ ImageNet ไว้อย่างสมบูรณ์\n• ประสิทธิภาพความเร็ว: Throughput 357 FPS ทำให้การทำ Dynamic Batching ใน DataLoader ไหลลื่น",
         COLOR_PRIMARY),

        ("ความคุ้มค่าและความเสี่ยงต่ำ (High ROI & Low Risk)",
         "• การใช้ Transfer Learning ลดเวลาในการ Train จากระดับสัปดาห์เหลือเพียงไม่กี่ชั่วโมงบน GPU\n• ไม่ต้องเก็บข้อมูลใหม่ทั้งหมด: อาศัย Pre-trained Feature Extractor จาก ImageNet ทำให้ผลลัพธ์โมเดลแม่นยำสูงและมีความเสถียร\n• Zero Data Leakage: การันตีว่าผลการประเมินจะไม่มี Overfitting หลอกตา",
         COLOR_ACCENT),

        ("คุณค่าเชิงการนำไปใช้จริง (Real-World Impact)",
         "• ตลาด Pet Industry เติบโตอย่างต่อเนื่อง การมีแอปพลิเคชันที่ระบุสายพันธุ์สัตว์เลี้ยงได้แม่นยำสร้าง Value ให้ทั้งผู้เลี้ยงและธุรกิจ\n• สนับสนุนงานด้านสวัสดิภาพสัตว์ (Shelter Adoption): ช่วยศูนย์ช่วยเหลือสัตว์จัดหมวดหมู่สัตว์จรจัดได้อย่างรวดเร็ว\n• สามารถต่อยอดเป็น AI API ให้บริการระบบภายนอกได้",
         RGBColor(124, 58, 237))
    ]

    card_w6 = Inches(3.72)
    card_h6 = Inches(4.8)
    gap6 = Inches(0.28)

    for i, (p_title, p_content, p_col) in enumerate(pillars):
        px = Inches(0.8) + i * (card_w6 + gap6)
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(1.8), card_w6, card_h6)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = p_col
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = p_title
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = p_col
        p1.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = p_content
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(12)

    # =========================================================================
    # SLIDE 7: Future Improvements (ข้อ 3 ในโจทย์ - สิ่งที่จะปรับปรุง)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "สิ่งที่จะปรับปรุงและพัฒนาต่อ หากมีเวลาหรือทรัพยากรเพิ่มเติม (Future Improvements)", "Future Improvements")

    areas = [
        ("1. การจัดการความสมดุลและคุณภาพของข้อมูล (Data Quality & Balance)",
         "• Targeted Data Collection: เก็บข้อมูลภาพเพิ่มเติมสำหรับคลาสที่มีภาพน้อยกว่า 100 ภาพ (เช่น Saint Bernard, Pinscher, Dachshund) ให้มีอย่างน้อย 300+ ภาพ/คลาส\n• Generative Data Augmentation: นำโมเดลสังเคราะห์ภาพ (Diffusion Models / ControlNet) มาช่วยสังเคราะห์สายพันธุ์ที่หายากอย่างสมจริง\n• Perceptual Deduplication: ใช้ pHash หรือ Embedding Cosine Similarity ตรวจจับและคัดกรองภาพที่ซ้ำซ้อนระดับมุมมองใกล้เคียงกัน"),

        ("2. การเพิ่มประสิทธิภาพการแบ่งข้อมูล (Data Splitting Optimization)",
         "• Re-Splitting สู่สัดส่วนมาตรฐาน: ปรับสัดส่วนจาก 92/5/3 ให้เป็น 80% Train / 10% Valid / 10% Test เพื่อให้ชุด Test มีขนาดใหญ่ขึ้น (2,800+ ภาพ) สร้างความมั่นใจในการวัดผล\n• Stratified K-Fold Cross Validation: จัดทำสคริปต์ทำ 5-Fold Cross Validation เพื่อการันตีว่าโมเดลมีความทนทาน (Robustness) ไม่ขึ้นกับสุ่มรอบใดรอบหนึ่ง"),

        ("3. การยกระดับสู่ระบบ MLOps เต็มรูปแบบ (Infrastructure & Automation)",
         "• Data Version Control (DVC): ติดตั้ง DVC ร่วมกับ Cloud Storage (S3 / GCS) สำหรับ Versioning ไฟล์ภาพ 28,000+ ภาพอย่างเป็นระบบ\n• Automated CI/CD Testing: ทำ GitHub Actions เพื่อรัน Unit Test และ Data Quality Verification ทุกครั้งที่มีการ Commit โค้ดใหม่\n• GPU-Accelerated Preprocessing: แปลง Preprocessing Pipeline เป็น Albumentations หรือ NVIDIA DALI เพื่อเร่ง Throughput จาก 357 FPS สู่ระดับ 2,000+ FPS บน GPU")
    ]

    card_y7 = Inches(1.8)
    c_h7 = Inches(1.55)
    gap7 = Inches(0.2)

    for i, (title, content) in enumerate(areas):
        card = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), card_y7 + i * (c_h7 + gap7), Inches(11.73), c_h7)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_BG
        card.line.color.rgb = COLOR_BORDER
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(12)
        p1.font.color.rgb = COLOR_PRIMARY
        p1.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = content
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT_DARK
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(4)

    # =========================================================================
    # SLIDE 8: Summary & Conclusion (Dark Theme)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    bg8 = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = COLOR_BG_DARK
    bg8.line.fill.background()

    # Title & Subtitle Box
    tbox8 = s8.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.33), Inches(1.5))
    tf8 = tbox8.text_frame
    tf8.word_wrap = True

    p8_1 = tf8.paragraphs[0]
    p8_1.text = "บทสรุปและก้าวต่อไป (Summary & Ready for Next Phase)"
    p8_1.font.size = Pt(32)
    p8_1.font.bold = True
    p8_1.font.color.rgb = COLOR_TEXT_LIGHT
    p8_1.font.name = "Segoe UI"

    p8_2 = tf8.add_paragraph()
    p8_2.text = "Data Pipeline เสร็จสมบูรณ์ตามหลักวิศวกรรม Machine Learning พร้อมก้าวสู่การเทรนและ Deploy โมเดล"
    p8_2.font.size = Pt(16)
    p8_2.font.color.rgb = COLOR_PRIMARY
    p8_2.font.name = "Segoe UI"
    p8_2.space_before = Pt(8)

    # 3 Summary Cards
    cards_summary = [
        ("1. Data Pipeline Complete", "ดาวน์โหลด, ตรวจสอบ, EDA, Preprocess (357 FPS) และแบ่งชุดข้อมูล Train/Val/Test 104 คลาสเรียบร้อย 100%"),
        ("2. Automation & Reproducibility", "ควบคุมผ่าน run_pipeline.bat รันจบในคลิกเดียว พร้อมเอกสารรายงาน Markdown และชาร์ตสถิติครบถ้วน"),
        ("3. Ready for SOTA Modeling", "รองรับการทำ Transfer Learning (EfficientNet/ConvNeXt) และพร้อมต่อยอดเป็น Real-time Edge Application")
    ]

    sc_w = Inches(3.64)
    sc_h = Inches(2.8)
    sc_gap = Inches(0.2)

    for i, (st, sc) in enumerate(cards_summary):
        px = Inches(1.0) + i * (sc_w + sc_gap)
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px, Inches(3.0), sc_w, sc_h)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_DARK
        card.line.color.rgb = RGBColor(51, 65, 85)
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = st
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = COLOR_ACCENT
        p1.font.name = "Segoe UI"

        p2 = tf.add_paragraph()
        p2.text = sc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_LIGHT
        p2.font.name = "Segoe UI"
        p2.space_before = Pt(12)

    # Footer Q&A
    ft = s8.shapes.add_textbox(Inches(1.0), Inches(6.2), Inches(11.33), Inches(0.6))
    p_ft = ft.text_frame.paragraphs[0]
    p_ft.text = "Repository: https://github.com/wayu0849-lang/-   |   Q&A Session & Thank You"
    p_ft.font.size = Pt(13)
    p_ft.font.color.rgb = RGBColor(148, 163, 184)
    p_ft.font.name = "Segoe UI"
    p_ft.alignment = PP_ALIGN.CENTER

    # Save presentations
    output_dir = Path("reports")
    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / "presentation_dogs_cats_pipeline.pptx"
    prs.save(str(out_path.resolve()))
    
    # Also save to root for easy access
    root_out = Path("presentation_dogs_cats_pipeline.pptx")
    prs.save(str(root_out.resolve()))

    print(f"[PPTX] Presentation slides successfully generated at: {out_path.resolve()} and {root_out.resolve()}")
    return out_path

if __name__ == "__main__":
    create_presentation()
